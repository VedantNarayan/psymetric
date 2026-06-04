from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement

# Image paths
img_student = "/Users/Vedant/.gemini/antigravity/brain/9813096c-b6ec-4d02-9585-e83ba47f3109/student_dashboard_v2_1773516117780.png"
img_school  = "/Users/Vedant/.gemini/antigravity/brain/9813096c-b6ec-4d02-9585-e83ba47f3109/school_dashboard_v2_1773516146695.png"
img_client  = "/Users/Vedant/.gemini/antigravity/brain/9813096c-b6ec-4d02-9585-e83ba47f3109/client_dashboard_v2_1773516173964.png"
img_logo    = "/Users/Vedant/Documents/PsyMetric/PsyMetric Logo (transparent).png"
img_mascot  = "/Users/Vedant/Documents/PsyMetric/Psy (Masscot).jpeg"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# UI Theme Colors (Clean Light Theme)
BG_COLOR = RGBColor(252, 252, 252)        # Very Light Minimal Background
CARD_BG = RGBColor(255, 255, 255)         # Pure White for Cards
ACCENT_BLUE = RGBColor(14, 165, 233)      # Modern Brand Blue 
ACCENT_LIGHT = RGBColor(241, 245, 249)    # Light dividers
BORDER_COLOR = RGBColor(226, 232, 240)    # Soft gray borders
TEXT_MAIN = RGBColor(30, 41, 59)          # Deep Gray/Charcoal for Titles
TEXT_MUTED = RGBColor(71, 85, 105)        # Medium Gray for body text

def add_transition(slide):
    """Add a smooth fade transition to the slide using raw OXML."""
    transition = OxmlElement('p:transition')
    transition.set('spd', 'med')
    fade = OxmlElement('p:fade')
    transition.append(fade)
    slide.element.insert(1, transition)

def animate_shape(slide, shape_id):
    """Adds a basic appearance animation to a shape."""
    timing = slide.element.xpath('./p:timing')
    if not timing:
        timing = OxmlElement('p:timing')
        slide.element.append(timing)
    else:
        timing = timing[0]

    tnLst = timing.xpath('./p:tnLst')
    if not tnLst:
        tnLst = OxmlElement('p:tnLst')
        timing.append(tnLst)
    else:
        tnLst = tnLst[0]

    par = OxmlElement('p:par')
    tnLst.append(par)
    
    cTn = OxmlElement('p:cTn')
    cTn.set('id', '1')
    cTn.set('nodeType', 'tmRoot')
    par.append(cTn)

    childTnLst = OxmlElement('p:childTnLst')
    cTn.append(childTnLst)
    
    seq = OxmlElement('p:seq')
    seq.set('concurrent', '1')
    seq.set('nextAc', 'seek')
    childTnLst.append(seq)

    # Simplified animation XML chunk (Appear)
    anim_xml = f"""
    <p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
        <p:cTn id="2" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" nodeType="withEffect">
            <p:stCondLst>
                <p:cond delay="0"/>
            </p:stCondLst>
            <p:childTnLst>
                <p:set>
                    <p:cBhvr>
                        <p:cTn id="3" dur="1" fill="hold">
                            <p:stCondLst>
                                <p:cond delay="0"/>
                            </p:stCondLst>
                        </p:cTn>
                        <p:tgtEl>
                            <p:spTgt spid="{shape_id}"/>
                        </p:tgtEl>
                        <p:attrNameLst>
                            <p:attrName>style.visibility</p:attrName>
                        </p:attrNameLst>
                    </p:cBhvr>
                    <p:to>
                        <p:strVal val="visible"/>
                    </p:to>
                </p:set>
            </p:childTnLst>
        </p:cTn>
    </p:par>
    """
    import lxml.etree as etree
    anim_elem = etree.fromstring(anim_xml)
    seq.append(anim_elem)

def set_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Add Logo
    slide.shapes.add_picture(img_logo, Inches(0.5), Inches(0.4), height=Inches(0.6))
    add_transition(slide)

def add_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.35), Inches(10), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Poppins'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    # Add a soft minimalist separator under the title
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.2), Inches(0.6), Inches(0.04))
    sep.fill.solid()
    sep.fill.fore_color.rgb = ACCENT_BLUE
    sep.line.fill.background() # No border

def create_card(slide, left, top, width, height, title, texts):
    # Card Background (Rounded Rectangle)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    
    # Very subtle modern border 
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(1.0)
    
    # Card Title
    title_box = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + 0.3), Inches(width - 0.6), Inches(0.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = 'Poppins'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    # Subtle separator line below title
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + 0.3), Inches(top + 0.85), Inches(width - 0.6), Inches(0.01))
    sep.fill.solid()
    sep.fill.fore_color.rgb = ACCENT_LIGHT
    sep.line.fill.background()
    
    # Card Text / Bullet points
    text_box = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + 1.0), Inches(width - 0.6), Inches(height - 1.2))
    tf = text_box.text_frame
    tf.word_wrap = True
    
    for i, t in enumerate(texts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = t
        p.font.name = 'Lato'
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(14)
        p.line_spacing = 1.3
        
    # Apply Entrance Animation to the text box
    try:
        animate_shape(slide, text_box.shape_id)
    except:
        pass 

# --- SLIDE 1: TITLE SLIDE ---
slide = prs.slides.add_slide(prs.slide_layouts[6]) # completely blank
set_background(slide)
slide.shapes.add_picture(img_logo, Inches(5.76), Inches(2.2), height=Inches(1.8))

title_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11.33), Inches(1))
p = title_box.text_frame.paragraphs[0]
p.text = "PsyMetric Labs PVT LTD"
p.font.name = 'Poppins'
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = TEXT_MAIN
p.alignment = PP_ALIGN.CENTER

sub_box = slide.shapes.add_textbox(Inches(1), Inches(5.3), Inches(11.33), Inches(1))
p = sub_box.text_frame.paragraphs[0]
p.text = "psymetric.online - The Future of Career Analysis for Students"
p.font.name = 'Lato'
p.font.size = Pt(20)
p.font.color.rgb = ACCENT_BLUE
p.alignment = PP_ALIGN.CENTER

# --- SLIDE 2: OVERVIEW & VISION ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_title(slide, "Overview & Vision")

create_card(slide, 1, 1.6, 5.3, 5, "Target Audience", [
    "Primarily targeting 8th/9th and 11th/12th Grade Students.",
    "Providing crucial guidance during their most pivotal educational milestones."
])

create_card(slide, 6.7, 1.6, 5.6, 5, "The AI-Driven Mechanism", [
    "Assessment Model: Initial pilot phase utilizes a fixed questionnaire. Transitions into a fully dynamic, AI-driven evaluation post-pilot.",
    "Technology Core: Questions are generated & individually tailored by our in-house trained AI model.",
    "Evaluation: Powered by a proprietary, research-backed Machine Learning algorithm.",
    "Ultimate Outcome: Highly accurate, data-driven career path guidance matched to the student's unique psychometric profile."
])

# --- DASHBOARD SLIDES OVERHAUL ---
def add_dashboard_slide(title_text, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, title_text)
    
    # Clean Light Background Container for Dashboard
    backplate = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.85), Inches(1.5), Inches(9.6), Inches(5.5))
    backplate.fill.solid()
    backplate.fill.fore_color.rgb = CARD_BG
    backplate.line.color.rgb = BORDER_COLOR
    backplate.line.width = Pt(1.0)
    
    # Image nicely padded inside
    img = slide.shapes.add_picture(img_path, Inches(1.95), Inches(1.6), width=Inches(9.4))
    try:
        animate_shape(slide, img.shape_id)
    except:
        pass

# --- SLIDE 3, 4, 5: DASHBOARDS ---
add_dashboard_slide("Student Dashboard (Personal Progress)", img_student)
add_dashboard_slide("School Dashboard (Aggregated Analytics)", img_school)
add_dashboard_slide("Client Dashboard (Direct Signup)", img_client)

# --- SLIDE 6: GAMIFICATION & ENGAGEMENT ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_title(slide, "Gamification & Student Engagement")

mascot = slide.shapes.add_picture(img_mascot, Inches(1), Inches(1.8), height=Inches(4.8))
try: animate_shape(slide, mascot.shape_id)
except: pass

create_card(slide, 5.5, 1.6, 6.8, 5, "Transforming the Assessment Experience", [
    "Game-Like Format: Questionnaires are designed to feel like an interactive game. We completely eliminate the stress, pressure, and fatigue of traditional diagnostic interviews.",
    "Reward & Progression: We implement a comprehensive Badge System. Students earn badges as part of the game to encourage completion and celebrate milestones.",
    "Mascot Integration: 'Psy' the owl serves as a friendly, relatable companion and guide throughout the entire user journey."
])

# --- SLIDE 7: POST-ASSESSMENT ECOSYSTEM ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_title(slide, "Post-Assessment Ecosystem")

create_card(slide, 1, 1.6, 11.3, 2.0, "Continuous Engagement", [
    "Our relationship with the user does not end after the assessment. We guide their next steps actively."
])
create_card(slide, 1, 3.9, 5.4, 2.8, "Experiential Learning", [
    "Tailored Summer Camps: Immersive, hands-on experiences customized for different student categories based on their assessment results.",
    "Interactive Workshops & Seminars: Targeted skill-building sessions hosted by specific industry experts."
])
create_card(slide, 6.9, 3.9, 5.4, 2.8, "Modern Media Access", [
    "Podcasts & Audio Series: Interviews with professionals tailored to students' matched career paths.",
    "Continuous Content Ecosystem: Keeping students deeply engaged with the PsyMetric brand."
])

# --- SLIDE 8: BUSINESS PLAN ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_title(slide, "Business Plan & Go-To-Market")

create_card(slide, 1, 1.6, 5.4, 5, "Phase 1: Pilot (0-6 Months)", [
    "Strategy: Provide the platform to 10 Schools completely FREE.",
    "Purpose: This phase uses the static questionnaire.",
    "Goal: Critical for collecting robust, real-world baseline data from thousands of students to train and refine our dynamic AI models before charging."
])
create_card(slide, 6.9, 1.6, 5.4, 5, "Phase 2: Validation", [
    "Strategy: Onboard 20 Schools (average of 200 students per school).",
    "Revenue Model: ₹100 per student.",
    "Expected Early Revenue: ₹4,00,000.",
    "Outcome: Prepares the business for immediate scaling and shows strong product-market fit."
])

# --- SLIDE 9: MARKET POTENTIAL ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_title(slide, "Expansion & Phenomenal Market Potential")

create_card(slide, 1, 1.6, 11.3, 5, "The Indian Landscape & The Gap", [
    "Current Value: India's career counseling market is estimated at over ₹5,000 Crore.",
    "Growth Trajectory: Experiencing a robust 15% CAGR, poised to double to ₹10,000 Crore by the end of the decade.",
    "The Massive Gap: There is a severe shortage of trained professionals. India currently has only ~10,000 counselors against a projected need of 14 Lakhs.",
    "Our Advantage: This creates an astronomical vacuum for accessible, AI-driven EdTech assessments. We scale instantly where human counselors cannot.",
    "Phase 3 Goal: Full regional dominance, starting with capturing the complete Bihar region, leading to massive 3-year profit margins."
])

# --- SLIDE 10: FINANCIALS ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_title(slide, "Financials: Expected Expenses")

create_card(slide, 1, 1.6, 5.4, 2.3, "Tech Infrastructure", ["Software development, cloud hosting, AI model training cycles, and ML algorithm maintenance."])
create_card(slide, 6.9, 1.6, 5.4, 2.3, "Operations", ["Day-to-day business management, platform scaling maintenance, and customer success support hubs."])
create_card(slide, 1, 4.2, 5.4, 2.6, "Sales & Collaborations", ["Strategic partnerships with educational institutions and hiring on-field B2B sales executives."])
create_card(slide, 6.9, 4.2, 5.4, 2.6, "Marketing & Giftings", ["Strategic gifts for faculties and school principals.", "Designed to build trust, expand network, and secure high-value B2B pilot programs and sales."])

prs.save('/Users/Vedant/Documents/PsyMetric/PsyMetric_Pitch.pptx')
print("Successfully generated beautifully polished Light Theme Presentation with Animations!")
